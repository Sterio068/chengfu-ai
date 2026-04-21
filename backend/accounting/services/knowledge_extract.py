"""
多格式抽字器 · V1.1-SPEC §E-2
=====================================
依副檔名路由 · 統一回 {type, content_preview, filename, size, modified_at, ...}
所有抽字失敗都不 raise · 回 {type: "error", error: str(e)}
讓 indexer 繼續處理下一個檔

本模組對外只暴露 extract(path) · 其他 helper 都 underscore
"""
import os
import logging
import pathlib
from datetime import datetime
from typing import Any

logger = logging.getLogger("chengfu.extract")

# lazy import · 讓測試環境不一定要裝滿所有 library
_fitz = None
_docx = None
_pptx = None
_openpyxl = None
_Image = None
_TAGS = None


def _lazy_fitz():
    global _fitz
    if _fitz is None:
        import fitz  # pymupdf
        _fitz = fitz
    return _fitz


def _lazy_docx():
    global _docx
    if _docx is None:
        from docx import Document
        _docx = Document
    return _docx


def _lazy_pptx():
    global _pptx
    if _pptx is None:
        from pptx import Presentation
        _pptx = Presentation
    return _pptx


def _lazy_openpyxl():
    global _openpyxl
    if _openpyxl is None:
        from openpyxl import load_workbook
        _openpyxl = load_workbook
    return _openpyxl


def _lazy_image():
    global _Image, _TAGS
    if _Image is None:
        from PIL import Image
        from PIL.ExifTags import TAGS
        _Image, _TAGS = Image, TAGS
    return _Image, _TAGS


# ------------------------------------------------------------
# Extractor · 每格式一個函式 · 失敗交由外層 extract() catch
# ------------------------------------------------------------
def _extract_pdf(path: str) -> dict:
    fitz = _lazy_fitz()
    doc = fitz.open(path)
    pages = []
    ocr_triggered = 0
    try:
        for page in doc:
            text = page.get_text("text").strip()
            # 若文字太少而頁面有圖片 · 嘗試 OCR 降級
            if len(text) < 120 and page.get_images():
                try:
                    tp = page.get_textpage_ocr(language="chi_tra+eng")
                    text = page.get_text(textpage=tp).strip()
                    ocr_triggered += 1
                except Exception:
                    # OCR 未安裝 tesseract 或其他錯 · 維持原 text
                    pass
            pages.append(text)
    finally:
        doc.close()
    return {
        "type": "pdf",
        "page_count": len(pages),
        "ocr_pages": ocr_triggered,
        "content_preview": ("\n\n".join(pages))[:2000],
    }


def _extract_docx(path: str) -> dict:
    Document = _lazy_docx()
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # 表格內容也抓(簡易)
    for table in doc.tables:
        for row in table.rows:
            row_txt = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_txt:
                paragraphs.append(row_txt)
    text = "\n".join(paragraphs)
    return {
        "type": "docx",
        "paragraph_count": len(paragraphs),
        "content_preview": text[:2000],
    }


def _extract_pptx(path: str) -> dict:
    Presentation = _lazy_pptx()
    prs = Presentation(path)
    slides_text = []
    for i, s in enumerate(prs.slides, start=1):
        parts = []
        for shape in s.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        if parts:
            slides_text.append(f"[投影片 {i}]\n" + "\n".join(parts))
    return {
        "type": "pptx",
        "slide_count": len(prs.slides),
        "content_preview": ("\n\n".join(slides_text))[:2000],
    }


def _extract_xlsx(path: str) -> dict:
    load_workbook = _lazy_openpyxl()
    wb = load_workbook(path, read_only=True, data_only=True)
    try:
        preview_parts = []
        # 最多前 3 張 sheet · 每張前 20 行
        for sheet_name in wb.sheetnames[:3]:
            ws = wb[sheet_name]
            rows = []
            for i, row in enumerate(ws.iter_rows(max_row=20, values_only=True)):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
                if i >= 19:
                    break
            if rows:
                preview_parts.append(f"[{sheet_name}]\n" + "\n".join(rows))
        return {
            "type": "xlsx",
            "sheet_count": len(wb.sheetnames),
            "content_preview": ("\n\n".join(preview_parts))[:2000],
        }
    finally:
        wb.close()


def _extract_image(path: str) -> dict:
    """設計圖只抽 metadata · 不抽內容(v2 再做 CLIP vision)"""
    Image, TAGS = _lazy_image()
    with Image.open(path) as img:
        width, height, fmt = img.width, img.height, img.format
        exif = {}
        try:
            raw = getattr(img, "_getexif", lambda: None)()
            if raw:
                for k, v in raw.items():
                    name = TAGS.get(k, str(k))
                    # 避免存大 binary · 截短
                    try:
                        exif[name] = str(v)[:100]
                    except Exception:
                        pass
        except Exception:
            pass
    return {
        "type": "image",
        "width": width,
        "height": height,
        "format": fmt,
        "exif": exif,
        "content_preview": f"圖片 {width}x{height} · {fmt}",
    }


def _extract_text(path: str) -> dict:
    """.txt / .md / .csv 直接 utf-8 讀"""
    # 最多讀 1MB · 避免超大 log 檔炸掉
    with open(path, "rb") as f:
        raw = f.read(1024 * 1024)
    for enc in ("utf-8", "utf-8-sig", "big5", "cp950", "latin-1"):
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = raw.decode("utf-8", errors="replace")
    return {
        "type": "text",
        "content_preview": text[:2000],
    }


# ------------------------------------------------------------
# 副檔名 → extractor 路由表
# ------------------------------------------------------------
EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx, ".xls": _extract_xlsx,
    ".jpg": _extract_image, ".jpeg": _extract_image,
    ".png": _extract_image, ".tiff": _extract_image, ".webp": _extract_image,
    ".txt": _extract_text, ".md": _extract_text, ".csv": _extract_text,
    # AI / PSD / AE 等格式無法讀內容 · fallback 走 unknown
}


def extract(path: str) -> dict:
    """對外主入口 · 統一回字典 · 失敗 type='error'"""
    ext = pathlib.Path(path).suffix.lower()
    result: dict[str, Any] = {}
    try:
        if ext in EXTRACTORS:
            result = EXTRACTORS[ext](path)
        else:
            result = {
                "type": "unknown",
                "content_preview": f"檔案 {pathlib.Path(path).name}（格式 {ext} 無法抽字）",
            }
    except Exception as e:
        logger.warning("[extract] %s · %s: %s", path, type(e).__name__, e)
        return {
            "path": path,
            "filename": pathlib.Path(path).name,
            "type": "error",
            "content_preview": "",
            "error": f"{type(e).__name__}: {e}"[:200],
        }
    # 補 file-level metadata
    try:
        stat = os.stat(path)
        result.update({
            "path": path,
            "filename": pathlib.Path(path).name,
            "size": stat.st_size,
            "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        })
    except Exception as e:
        result["stat_error"] = str(e)
    return result
